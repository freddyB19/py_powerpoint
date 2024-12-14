import os

from pptx import Presentation

BASE = os.path.dirname(os.path.abspath(__file__))
FILENAME = 'trial.pptx'


def python_pptx():
	file = Presentation(os.path.join(BASE, 'file', FILENAME))

	for slide  in file.slides:
		for shape in slide.shapes:
			if not shape.has_text_frame:
				continue
			
			parrafo = shape.text_frame.paragraphs[0]
			nuevo_texto = parrafo.runs[0]
			nuevo_texto.text = "Hola"


	file.save(os.path.join(BASE, 'file', FILENAME))
	
	


def main():
	python_pptx()

if __name__ == '__main__':
	main()